# 🧩 Required Libraries
import os
import fitz  # PyMuPDF
import docx
import requests
from pptx import Presentation
from elevenlabs import ElevenLabs
from openai import OpenAI

# 🔐 API KEYS
eleven_key = ElevenLabs(
    api_key=userdata.get('ELEVENLABS_API_KEY'),
)
openrouter_key = "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"

client = OpenAI(
    base_url="https://openrouter.ai/api/v1",
    api_key=openrouter_key,
)

# 🧠 Simplify content using OpenRouter (Qwen)
conversation_history = []

def simplify_prompt(user_text):
    # Add user question to history
    conversation_history.append({"role": "user", "content": user_text})

    completion = client.chat.completions.create(
        model="qwen/qwen3-30b-a3b:free",
        extra_headers={
            "HTTP-Referer": "https://neurobridge.ai",
            "X-Title": "NeuroBridge",
        },
        messages=conversation_history
    )

    response_text = completion.choices[0].message.content.strip()
    conversation_history.append({"role": "assistant", "content": response_text})
    return response_text

# 🎤 Text-to-Speech with ElevenLabs
from IPython.display import Audio, display
from io import BytesIO

def speak_with_tone(text):
    audio = eleven_key.text_to_speech.convert(
        voice_id="W7iR5kTNHozpIl2Jqq15",
        output_format="mp3_44100_128",
        text=text,
        model_id="eleven_multilingual_v2",
    )

    # # Join generator into BytesIO buffer
    # audio_bytes = BytesIO(b"".join(audio))
    # audio_bytes.seek(0)

    # # Display and autoplay audio in Colab
    # display(Audio(audio_bytes.read(), autoplay=True))

# 💬 Text-based feedback
def get_user_feedback():
    return input("👤 Did you understand? (yes/no): ").strip().lower()

# 📄 File Text Extractors
def extract_text_from_pdf(file_path):
    text = ''
    for page in fitz.open(file_path):
        text += page.get_text()
    return text.strip()

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    return ' '.join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

def extract_text_from_docx(file_path):
    return '\n'.join([p.text for p in docx.Document(file_path).paragraphs if p.text.strip()])

# 🔁 File Loader
def load_text_from_file(file_path):
    if file_path.endswith(".pdf"):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith(".pptx"):
        return extract_text_from_pptx(file_path)
    elif file_path.endswith(".docx"):
        return extract_text_from_docx(file_path)
    else:
        raise ValueError("Unsupported file format.")

import requests

def generate_flux_image(prompt):
    api_key = "82ece252-62b4-4e4e-9f57-d7a495b0af62:307ec33823f5bba87ac0a52dd71a6962"
    url = "https://fal.run/fal-ai/flux/schnell"
    
    headers = {
        "Authorization": f"Key {api_key}",
        "Content-Type": "application/json"
    }

    payload = {
            "prompt": prompt
    }

    response = requests.post(url, json=payload, headers=headers)

    if response.status_code == 200:
        print("✅ Image generated!")
        print(response.json())
        return response.json().get("image", "No image returned")
    else:
        print("❌ Error generating image:", response.status_code, response.text)
        return None

# 📊 Visual Explanation from OpenRouter (InternVL model)
def generate_visual_explanation(topic_text):
    visual_prompt = f"Create a detailed visual illustration or flowchart to explain this concept to a child with learning difficulties:\n\n{topic_text}"

    # Generate prompt from Qwen
    prompt_response = client.chat.completions.create(
        model="qwen/qwen3-30b-a3b:free",
        extra_headers={
            "HTTP-Referer": "https://neurobridge.ai",
            "X-Title": "NeuroBridge",
        },
        messages=[{"role": "user", "content": visual_prompt}]
    )

    generated_prompt = prompt_response.choices[0].message.content.strip()
    print("\n🎨 Visual Prompt for Flux:\n", generated_prompt)

# 🔁 Full Conversation Loop
def neurobridge_conversation_loop(topic_text):
    print("🧠 Starting NeuroBridge...\n")

    # Start with a system message
    conversation_history.clear()
    conversation_history.append({
        "role": "system",
        "content": "You are a friendly, playful, engaging tutor for K-12 students with dyslexia, ADHD, or autism. Always explain simply and end with: 'Did that make sense?'"
    })

    # Begin conversation
    explanation = simplify_prompt(f"Please explain this topic:\n\n{topic_text}")
    print("🤖 NeuroBuddy says:\n", explanation)
    speak_with_tone(explanation)

    while True:
        feedback = get_user_feedback()
        if feedback in ["yes", "got it", "understood", "ok", "okay"]:
            speak_with_tone("Yay! You're doing amazing. Let’s move forward!")
            break
        else:
            followup = input("🔄 Which part confused you? Please type a keyword or sentence: ")
            updated = simplify_prompt(f"Can you explain more clearly about: {followup}")
            print("\n🤖 NeuroBuddy says:\n", updated)
            speak_with_tone(updated)

    # Visual Output
    generate_visual_explanation(topic_text)

# 📥 Upload and Run (Colab)
from google.colab import files
uploaded = files.upload()
file_path = list(uploaded.keys())[0]
topic_text = load_text_from_file(file_path)

# 🚀 Start Session
neurobridge_conversation_loop(topic_text)
